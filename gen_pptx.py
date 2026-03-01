"""
gen_pptx.py — Assemble slides into a 16:9 PPTX presentation.

Supports two modes:
  1. Image-only (legacy): --images slide_1.png slide_2.png --output out.pptx
  2. Dual-layer (new):    --deck deck.json --imgdir slides_tmp/ --output out.pptx

Dual-layer slides have:
  - Layer 1: rough.js background PNG (decorations only)
  - Layer 2: Native python-pptx TextBox / Image elements (editable)

Slides with render_mode="image_only" or no elements[] fall back to full-bleed PNG.
"""

import argparse
import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 16:9 widescreen dimensions
SLIDE_WIDTH = Inches(13.333)   # 13.333 inches = standard 16:9
SLIDE_HEIGHT = Inches(7.5)     # 7.5 inches

# Coordinate conversion: 1920px canvas → slide EMU width
# 12192000 EMU / 1920 px ≈ 6350 EMU/px
PX_TO_EMU = SLIDE_WIDTH / 1920  # Emu per pixel (float division of Emu values)

# Default font mapping: HTML font category → PowerPoint font
DEFAULT_FONT_MAP = {
    "heading": "Ink Free",
    "body": "Ink Free",
    "annotation": "Ink Free",
    "cjk": "Microsoft JhengHei",
}

# Alignment mapping
ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def px(val: int) -> int:
    """Convert pixel value to EMU."""
    return int(val * PX_TO_EMU)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _has_cjk(text: str) -> bool:
    """Check if text contains CJK characters."""
    for ch in text:
        cp = ord(ch)
        if (0x4E00 <= cp <= 0x9FFF or    # CJK Unified
            0x3400 <= cp <= 0x4DBF or    # CJK Extension A
            0x3000 <= cp <= 0x303F or    # CJK Symbols
            0x3040 <= cp <= 0x309F or    # Hiragana
            0x30A0 <= cp <= 0x30FF or    # Katakana
            0xAC00 <= cp <= 0xD7AF):     # Hangul
            return True
    return False


def _set_cjk_font(run, cjk_font: str) -> None:
    """Set East Asian font on a run via XML for CJK/Latin font switching."""
    r_elem = run._r
    rPr = r_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    if rPr is None:
        rPr = etree.SubElement(
            r_elem,
            '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr',
        )
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        ea = etree.SubElement(
            rPr,
            '{http://schemas.openxmlformats.org/drawingml/2006/main}ea',
        )
    ea.set('typeface', cjk_font)


def _set_body_anchor(tf, anchor: str) -> None:
    """Set vertical alignment on a text frame via bodyPr anchor attribute.

    Args:
        tf: TextFrame object.
        anchor: "top", "middle"/"ctr", or "bottom".
    """
    anchor_map = {"top": "t", "middle": "ctr", "ctr": "ctr", "bottom": "b"}
    anchor_val = anchor_map.get(anchor, "t")
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    body_pr = tf._txBody.find(f'{ns}bodyPr')
    if body_pr is not None:
        body_pr.set('anchor', anchor_val)


def _set_fill_opacity(shape, opacity: float) -> None:
    """Set fill transparency on a shape via XML alpha attribute.

    Args:
        shape: A python-pptx Shape with solid fill already applied.
        opacity: 0.0 (fully transparent) to 1.0 (fully opaque).
    """
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    # spPr may be under presentationml or drawingml namespace
    sp_pr = None
    for child in shape._element:
        if child.tag.endswith('}spPr'):
            sp_pr = child
            break
    if sp_pr is None:
        return
    solid_fill = sp_pr.find(f'{ns}solidFill')
    if solid_fill is None:
        return
    # Find the color element (srgbClr or similar)
    for color_elem in solid_fill:
        alpha = color_elem.find(f'{ns}alpha')
        if alpha is None:
            alpha = etree.SubElement(color_elem, f'{ns}alpha')
        # OOXML alpha: 0 = fully transparent, 100000 = fully opaque
        alpha.set('val', str(int(opacity * 100000)))


def _add_textbox(slide, elem: dict, font_map: dict) -> None:
    """Add a native TextBox element to the slide."""
    left = px(elem.get("x", 0))
    top = px(elem.get("y", 0))
    width = px(elem.get("w", 400))
    height = px(elem.get("h", 100))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Vertical alignment
    v_align = elem.get("valign", "top")
    _set_body_anchor(tf, v_align)

    paragraphs = elem.get("paragraphs", [])
    cjk_font = font_map.get("cjk", DEFAULT_FONT_MAP["cjk"])

    for p_idx, para_data in enumerate(paragraphs):
        if p_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Paragraph-level formatting
        align = para_data.get("align", "left")
        p.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)

        # Spacing
        space_before = para_data.get("space_before")
        if space_before is not None:
            p.space_before = Pt(space_before)
        space_after = para_data.get("space_after")
        if space_after is not None:
            p.space_after = Pt(space_after)

        # Bullet support
        if para_data.get("bullet"):
            p.level = para_data.get("level", 0)

        # Handle runs: either a single text string or explicit runs[]
        runs_data = para_data.get("runs")
        if runs_data is None:
            # Single text paragraph — treat as one run
            runs_data = [{
                "text": para_data.get("text", ""),
                "font": para_data.get("font", "body"),
                "size": para_data.get("size", 28),
                "color": para_data.get("color"),
                "bold": para_data.get("bold", False),
                "italic": para_data.get("italic", False),
            }]

        for run_data in runs_data:
            run = p.add_run()
            run.text = run_data.get("text", "")

            font_cat = run_data.get("font", "body")
            font_name = font_map.get(font_cat, font_map.get("body", DEFAULT_FONT_MAP["body"]))
            run.font.name = font_name

            size = run_data.get("size", 28)
            run.font.size = Pt(size)

            color = run_data.get("color")
            if color:
                run.font.color.rgb = hex_to_rgb(color)

            if run_data.get("bold"):
                run.font.bold = True
            if run_data.get("italic"):
                run.font.italic = True

            # Set CJK font on every run for automatic CJK/Latin switching
            _set_cjk_font(run, cjk_font)


def _add_image_element(slide, elem: dict, style: dict = None) -> None:
    """Add a native Image element or a placeholder rectangle to the slide.

    Placeholder supports:
      - bg_color / bg_opacity: tinted background fill (intentional, not broken)
      - border_color: accent-colored border (falls back to style accent or gray)
      - icon: emoji character rendered above label text
      - label_color: custom label text color
    """
    left = px(elem.get("x", 0))
    top = px(elem.get("y", 0))
    width = px(elem.get("w", 400))
    height = px(elem.get("h", 300))

    src = elem.get("src")
    if src and Path(src).exists():
        slide.shapes.add_picture(src, left, top, width, height)
    else:
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )

        # Tinted background fill
        bg_color = elem.get("bg_color")
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(bg_color)
            bg_opacity = elem.get("bg_opacity", 0.3)
            _set_fill_opacity(shape, bg_opacity)
        else:
            shape.fill.background()  # No fill (backwards compat)

        # Border color: element-level > style accent > default gray
        border_color = elem.get("border_color")
        if not border_color and style:
            # Try accent palette from style
            accents = style.get("accents")
            if accents:
                border_color = accents[0]
        if not border_color:
            border_color = "#adb5bd"
        shape.line.color.rgb = hex_to_rgb(border_color)
        shape.line.width = Pt(1.5)
        shape.line.dash_style = 4  # dashed

        # Label with vertical centering
        tf = shape.text_frame
        tf.word_wrap = True
        _set_body_anchor(tf, "ctr")

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # Optional icon above label
        icon = elem.get("icon")
        if icon:
            icon_run = p.add_run()
            icon_run.text = icon + "\n"
            icon_run.font.size = Pt(36)
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = elem.get("label", "[Insert image]")
        run.font.size = Pt(18)
        label_color = elem.get("label_color", "#868e96")
        run.font.color.rgb = hex_to_rgb(label_color)
        run.font.name = "Ink Free"


def _add_fullbleed_image(slide, img_path: str) -> None:
    """Add a full-bleed image covering the entire slide."""
    slide.shapes.add_picture(
        img_path,
        left=Emu(0), top=Emu(0),
        width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
    )


def create_dual_pptx(deck: dict, imgdir: Path, output_path: str) -> None:
    """Create a PPTX with dual-layer slides from a deck JSON."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    style = deck.get("style", {})
    font_map = {**DEFAULT_FONT_MAP, **style.get("font_map", {})}

    slides_data = deck.get("slides", [])
    slide_count = 0

    for i, slide_data in enumerate(slides_data):
        slide_id = slide_data.get("id", i + 1)
        render_mode = slide_data.get("render_mode", "dual")
        elements = slide_data.get("elements", [])

        slide = prs.slides.add_slide(blank_layout)

        if render_mode == "image_only" or not elements:
            # Full-bleed PNG (current behavior)
            full_png = imgdir / f"slide_{slide_id:02d}.png"
            if full_png.exists():
                _add_fullbleed_image(slide, str(full_png))
            else:
                print(f"WARNING: Image not found for slide {slide_id}: {full_png}",
                      file=sys.stderr)
                continue
        else:
            # Dual-layer: bg PNG at z=0, then native elements on top
            bg_png = imgdir / f"slide_{slide_id:02d}_bg.png"
            if bg_png.exists():
                _add_fullbleed_image(slide, str(bg_png))
            else:
                # Fall back to full PNG if bg not available
                full_png = imgdir / f"slide_{slide_id:02d}.png"
                if full_png.exists():
                    print(f"WARNING: No bg image for slide {slide_id}, "
                          f"using full render", file=sys.stderr)
                    _add_fullbleed_image(slide, str(full_png))

            # Add native elements on top
            for elem in elements:
                elem_type = elem.get("type", "")
                if elem_type == "textbox":
                    _add_textbox(slide, elem, font_map)
                elif elem_type == "image":
                    _add_image_element(slide, elem, style)

        slide_count += 1

    if slide_count == 0:
        print("ERROR: No slides created. PPTX not generated.", file=sys.stderr)
        sys.exit(1)

    prs.save(output_path)
    print(f"Created: {output_path} ({slide_count} slides)")


def create_pptx(image_paths: list[str], output_path: str, title: str = "Presentation") -> None:
    """Create a PPTX from a list of PNG images, one image per slide (legacy mode)."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    for img_path in image_paths:
        p = Path(img_path)
        if not p.exists():
            print(f"WARNING: Image not found, skipping: {img_path}", file=sys.stderr)
            continue

        slide = prs.slides.add_slide(blank_layout)
        _add_fullbleed_image(slide, str(p))

    if len(prs.slides) == 0:
        print("ERROR: No valid images found. PPTX not created.", file=sys.stderr)
        sys.exit(1)

    prs.save(output_path)
    print(f"Created: {output_path} ({len(prs.slides)} slides)")


def main() -> None:
    parser = argparse.ArgumentParser(description="Assemble slides into PPTX")

    # Dual-layer mode args
    parser.add_argument("--deck", help="Path to JSON deck file (dual-layer mode)")
    parser.add_argument("--imgdir", help="Directory containing rendered PNGs")

    # Legacy mode args
    parser.add_argument(
        "--images", nargs="+",
        help="Ordered list of PNG image paths (legacy image-only mode)",
    )

    # Common args
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    parser.add_argument("--title", default="Presentation", help="Presentation title")

    args = parser.parse_args()

    if args.deck:
        # Dual-layer mode
        deck_path = Path(args.deck)
        if not deck_path.exists():
            print(f"ERROR: Deck file not found: {deck_path}", file=sys.stderr)
            sys.exit(1)

        with open(deck_path, "r", encoding="utf-8") as f:
            deck = json.load(f)

        imgdir = Path(args.imgdir) if args.imgdir else deck_path.parent / "slides_tmp"
        create_dual_pptx(deck, imgdir, args.output)

    elif args.images:
        # Legacy image-only mode
        create_pptx(args.images, args.output, args.title)

    else:
        print("ERROR: Provide either --deck (dual-layer) or --images (legacy).",
              file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
